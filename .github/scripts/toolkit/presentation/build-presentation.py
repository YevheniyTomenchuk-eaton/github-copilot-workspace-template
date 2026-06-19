#!/usr/bin/env python3
"""Build a polished, on-brand PowerPoint .pptx deck from a JSON content spec.

Usage:
    python build-presentation.py <spec.json> <output.pptx>

The deck is 16:9 widescreen with a branded title slide, optional section
dividers, and content slides that carry a colored title bar, an accent
underline, and a footer with the deck name and slide number.

Spec shape (see the office-documents skill for the full reference):
    {
      "title": "Q3 Business Review",
      "subtitle": "Sales & Operations",
      "author": "Dana Ruiz",
      "footer": "Q3 Business Review",
      "slides": [
        {"layout": "section", "title": "Where we are"},
        {"layout": "bullets", "kicker": "Overview", "title": "Three wins",
         "bullets": ["Revenue up 18%", "Two logos signed", "Churn down"]},
        {"layout": "two-column", "title": "Plan vs. actual",
         "left": ["Target: $4.0M"], "right": ["Actual: $4.3M"]},
        {"layout": "quote", "quote": "Best onboarding ever.",
         "attribution": "Pilot customer"}
      ]
    }

Emits machine-readable KEY=value lines on success:
    OUTPUT=<absolute path to the .pptx>
    SLIDES=<number of content slides>

Requires: python-pptx  (pip install python-pptx)
"""

import json
import os
import sys

PRIMARY = "1F4E79"
ACCENT = "00B3A4"
INK = "1A2233"
MUTED = "5B6B7F"
BAND = "EEF3F9"
WHITE = "FFFFFF"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _rgb(hex_color):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(hex_color)


def _add_box(slide, left, top, width, height, fill_hex=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.shadow.inherit = False
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_hex)
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def _text(slide, left, top, width, height, text, size, color_hex,
          bold=False, align="left", font="Calibri"):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color_hex)
    return box


def _bullets(slide, left, top, width, height, items, size=18):
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(10)
        marker = paragraph.add_run()
        marker.text = "\u25aa  "
        marker.font.size = Pt(size)
        marker.font.bold = True
        marker.font.name = "Calibri"
        marker.font.color.rgb = _rgb(ACCENT)
        run = paragraph.add_run()
        run.text = str(item)
        run.font.size = Pt(size)
        run.font.name = "Calibri"
        run.font.color.rgb = _rgb(INK)
    return box


def _blank_slide(presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _title_slide(presentation, spec):
    slide = _blank_slide(presentation)
    _add_box(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill_hex=PRIMARY)
    _add_box(slide, 0.9, 4.05, 2.2, 0.09, fill_hex=ACCENT)
    _text(slide, 0.9, 2.4, 11.5, 1.6, spec.get("title", "Untitled"),
          44, WHITE, bold=False, font="Calibri Light")
    if spec.get("subtitle"):
        _text(slide, 0.9, 4.3, 11.5, 0.8, spec["subtitle"], 22, "D7E3F2")
    footer_bits = [bit for bit in (spec.get("author"), spec.get("date")) if bit]
    if footer_bits:
        _text(slide, 0.9, 6.7, 11.5, 0.5, "   |   ".join(footer_bits), 12, "AFC2DB")


def _section_slide(presentation, slide_spec):
    slide = _blank_slide(presentation)
    _add_box(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill_hex=INK)
    _add_box(slide, 0.9, 3.55, 1.6, 0.08, fill_hex=ACCENT)
    _text(slide, 0.9, 3.0, 11.5, 1.2, slide_spec.get("title", ""),
          34, WHITE, bold=True)


def _content_header(slide, slide_spec):
    _add_box(slide, 0, 0, SLIDE_W_IN, 1.35, fill_hex=PRIMARY)
    _add_box(slide, 0.9, 1.28, SLIDE_W_IN - 1.8, 0.06, fill_hex=ACCENT)
    if slide_spec.get("kicker"):
        _text(slide, 0.92, 0.28, 11.5, 0.4,
              str(slide_spec["kicker"]).upper(), 12, "AFC2DB", bold=True)
        title_top = 0.62
    else:
        title_top = 0.42
    _text(slide, 0.9, title_top, 11.5, 0.8, slide_spec.get("title", ""),
          26, WHITE, bold=True)


def _footer(slide, footer_text, number):
    _text(slide, 0.9, 7.0, 9.0, 0.4, footer_text, 11, MUTED)
    _text(slide, 11.6, 7.0, 0.9, 0.4, str(number), 11, MUTED, align="right")


def _bullets_slide(presentation, slide_spec, footer_text, number):
    slide = _blank_slide(presentation)
    _content_header(slide, slide_spec)
    _bullets(slide, 0.9, 1.85, 11.5, 4.8, slide_spec.get("bullets", []))
    _footer(slide, footer_text, number)


def _two_column_slide(presentation, slide_spec, footer_text, number):
    slide = _blank_slide(presentation)
    _content_header(slide, slide_spec)
    _add_box(slide, 0.9, 1.85, 5.55, 4.6, fill_hex=BAND)
    _add_box(slide, 6.9, 1.85, 5.55, 4.6, fill_hex=BAND)
    _bullets(slide, 1.2, 2.15, 5.0, 4.0, slide_spec.get("left", []))
    _bullets(slide, 7.2, 2.15, 5.0, 4.0, slide_spec.get("right", []))
    _footer(slide, footer_text, number)


def _quote_slide(presentation, slide_spec, footer_text, number):
    slide = _blank_slide(presentation)
    _add_box(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, fill_hex=BAND)
    _add_box(slide, 0.9, 2.2, 0.12, 2.6, fill_hex=ACCENT)
    _text(slide, 1.4, 2.4, 10.6, 2.4,
          "\u201c" + str(slide_spec.get("quote", "")) + "\u201d",
          30, INK, bold=False, font="Calibri Light")
    if slide_spec.get("attribution"):
        _text(slide, 1.4, 4.9, 10.6, 0.6,
              "\u2014 " + str(slide_spec["attribution"]), 16, MUTED)
    _footer(slide, footer_text, number)


def build(spec_path, output_path):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("ERROR=python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    with open(spec_path, encoding="utf-8") as spec_file:
        spec = json.load(spec_file)

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W_IN)
    presentation.slide_height = Inches(SLIDE_H_IN)

    _title_slide(presentation, spec)

    footer_text = spec.get("footer") or spec.get("title", "")
    slides = spec.get("slides", [])
    number = 0
    for slide_spec in slides:
        layout = slide_spec.get("layout", "bullets")
        if layout == "section":
            _section_slide(presentation, slide_spec)
            continue
        number += 1
        if layout == "two-column":
            _two_column_slide(presentation, slide_spec, footer_text, number)
        elif layout == "quote":
            _quote_slide(presentation, slide_spec, footer_text, number)
        else:
            _bullets_slide(presentation, slide_spec, footer_text, number)

    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    presentation.save(output_path)

    print(f"OUTPUT={output_path}")
    print(f"SLIDES={len(slides)}")


def main():
    if len(sys.argv) != 3:
        print("ERROR=usage: python build-presentation.py <spec.json> <output.pptx>")
        sys.exit(1)
    build(sys.argv[1], sys.argv[2])


if __name__ == "__main__":
    main()
